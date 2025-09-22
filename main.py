import re
from functools import partial
from multiprocessing import cpu_count
from typing import Optional

import requests
from bs4 import BeautifulSoup
from bs4.element import Tag
from pptx import Presentation
from tqdm.contrib.concurrent import process_map

# If None, export all
PSALM: Optional[int] = None

superscript_lut = dict(zip((map(str, range(10))), "⁰¹²³⁴⁵⁶⁷⁸⁹"))  # From https://lingojam.com/TinyTextGenerator
superscript_lut.update({"a": "ᵃ", "b": "ᵇ"})


def all_after_verse_number(verse_element: Tag) -> [str, str]:
    vn_span = verse_element.find('span', class_='vn')
    # Extract all the text that comes after the 'vn' span
    text_after_vn = ''
    for element in vn_span.next_siblings:
        if isinstance(element, str):
            text_after_vn += element
        elif isinstance(element, Tag):
            text_after_vn += element.get_text()
    return text_after_vn.strip(), vn_span.get_text()


def grab_psalm(psalm):
    # URL of the website to scrape
    url = f"https://bibel.github.io/EUe/ot/Ps_{psalm}.html"

    # Send a GET request to the URL
    response = requests.get(url)

    # Check if the request was successful
    if response.status_code == 200:
        # Parse the HTML content
        soup = BeautifulSoup(response.content, 'html.parser')

        # Find the element with class "biblehtmlcontent verses"
        content_div = soup.find('div', class_='biblehtmlcontent verses')

        if content_div:
            # Find all elements with class "v" within the content_div
            verses = content_div.find_all('div', class_='v')

            # Extract the verses in their individual lines and their number
            verse_objects = []
            search_for_verse_purpose = True

            ps119_letter = None
            for verse in verses:
                verse_text_raw, verse_number = all_after_verse_number(verse)
                # Preprocessing: Remove trailing [sela], "/" or footnote number
                for _ in range(3):
                    # Repeat multiple times because footnote or sela could occur together
                    verse_text_raw = re.sub(r"(/$)|(\d+$)|(\[Sela])|(\[Zwischenspiel. Sela])$", "",
                                            verse_text_raw).strip()
                # Replace break dashes (they have no character on either side)
                verse_text_raw = re.sub(r"(?<![a-zA-Z])-(?![a-zA-Z])", "—", verse_text_raw)
                verse_text = verse_text_raw

                # Handle parentheses and hebrew letter in ps119
                if psalm == 119 and int(verse_number) != 1:
                    if ps119_letter is not None:
                        # There was a letter on the end of the last verse --> append it at the beginning here
                        verse_text = f"({ps119_letter}) {verse_text}"

                    if (match := re.findall(r"\((\w+)\)$", verse_text)) != []:
                        ps119_letter = match[0]
                        verse_text = verse_text.replace(f"({ps119_letter})", "").strip()
                    else:
                        ps119_letter = None

                # Remove any parentheses except the leading ones in ps 119
                verse_text = re.sub(r'(?<!^)\([^()]*\)', lambda m: m.group(0)[1:-1], verse_text)

                # verse text can either:
                # 1. already contain the correct first sentence (e.g. 107)
                # 2. contain only the psalms purpose in [ ] brackets (e.g. 85)
                # 3. contain part of the psalms purpose, spread over multiple verses (e.g. 51)
                # 4. contain the psalms purpose and the correct first sentence (e.g. 86)
                if search_for_verse_purpose and verse_text.startswith("["):
                    # Psalms purpose starting
                    if verse_text.endswith("]"):
                        # Psalms purpose spreads the whole current verse -> skip
                        search_for_verse_purpose = False
                        continue
                    elif "]" not in verse_text:
                        # Psalms purpose spreads more than the current verse -> skip at least this one
                        continue
                    else:
                        # Psalms purpose spreads part of the current verse -> cut it out
                        verse_text = re.sub(r"\[.+?]", "", verse_text).strip()
                        search_for_verse_purpose = False
                elif search_for_verse_purpose and "]" in verse_text:
                    if verse_text.endswith("]"):
                        # This verse ends completely with the psalms purpose -> skip it entirely
                        search_for_verse_purpose = False
                        continue
                    else:
                        # Psalms purpose continued here since last verse -> only cut the first part
                        verse_text = re.sub(r"^.+?]", "", verse_text).strip()
                        search_for_verse_purpose = False
                # Once the verse purpose stuff is out of the way, all other occurrences of brackets mean more text
                # that is just not handed down as reliably --> include it but remove the brackets.
                verse_text = verse_text.replace("[", "").replace("]", "")

                # Within-verse linebreaks are marked with forward slashes
                verse_lines = [v.strip() for v in verse_text.split("/")]

                # At the first line of a verse, add the verse number in superscript
                verse_lines[0] = "".join(superscript_lut[num] for num in verse_number) + " " + verse_lines[0]
                verse_objects.append(verse_lines)

            # At the last line of the last verse, append a Halleluja
            # verse_objects[-1] += " — Halleluja!"
            # verse_objects.append(["— Halleluja!"])  # Append it as a standalone line
            return verse_objects
        else:
            print("The specified class 'biblehtmlcontent verses' was not found.")
            raise ValueError()
    else:
        print(f"Failed to retrieve the webpage. Status code: {response.status_code}")
        raise ConnectionError()


class PsalmWriter:
    top_margin = 0.05  # Inches

    def __init__(self, psalm_number: int, prs: Presentation, body_font_size, line_spacing_factor, space_after, wrap):
        self.psalm_number = psalm_number
        self.prs = prs
        self.body_font_size = body_font_size
        self.line_spacing_factor = line_spacing_factor
        self.space_after = space_after
        self.warp = wrap
        self.width_inch = prs.slide_width.inches
        self.height_inch = prs.slide_height.inches
        self.current_text_body_height = 0

        self.verses = grab_psalm(self.psalm_number)

    def add_slide(self, title_slide: bool):
        # Add a slide with a title bar and a body text element
        slide_layout = self.prs.slide_layouts[1 - int(title_slide)]
        return self.prs.slides.add_slide(slide_layout)

    @staticmethod
    def get_text_frame(slide, is_title_slide=False):
        # Get the body text element:
        # This is the second element in the title slide master, otherwise the first element
        text_frame = slide.placeholders[int(is_title_slide)].text_frame

        return text_frame, text_frame.paragraphs[0]

    @staticmethod
    def fill_title(slide, title):
        # Fill in the title text element
        title_text_frame = slide.placeholders[0].text_frame
        title_text_frame.text = title

    def point_to_inch(self, line_spacing_factor=1.0, after=0.0, point_to_pixel=1.33, dpi=96):
        pixels = (self.body_font_size * line_spacing_factor + after) * point_to_pixel
        return pixels / dpi

    def write_psalm(self):
        slide = self.add_slide(title_slide=True)
        self.fill_title(slide, f"Psalm {self.psalm_number}")
        text_frame, p = self.get_text_frame(slide, is_title_slide=True)
        self.current_text_body_height = 0

        for i, paragraph_lines in enumerate(self.verses):
            # Calculate if verse still fits on slide
            num_lines = len(paragraph_lines)
            # Hotfix for line wrapping for very long lines (e.g. Ps 59,14 or Ps 145,13)
            wrapped_lines = sum([len(paragraph_line) > self.warp for paragraph_line in paragraph_lines])
            num_lines += wrapped_lines
            paragraph_height = (num_lines - 1) * self.point_to_inch(line_spacing_factor=self.line_spacing_factor)
            paragraph_height += self.point_to_inch(line_spacing_factor=1.0, after=self.space_after)
            if text_frame._parent.top.inches + self.top_margin + self.current_text_body_height + paragraph_height > self.height_inch:
                # Add a slide with only a body text element
                slide = self.add_slide(title_slide=False)
                self.current_text_body_height = 0
                text_frame, p = self.get_text_frame(slide, is_title_slide=True)
            elif i > 0:
                p = text_frame.add_paragraph()

            self.current_text_body_height += paragraph_height

            # Add text to paragraph
            for line_idx, line in enumerate(paragraph_lines):
                run = p.add_run()
                run.text = line
                if line_idx < len(paragraph_lines) - 1:
                    # Join the lines with soft line breaks:4
                    # See https://github.com/scanny/python-pptx/issues/322#issuecomment-339607317
                    p._p.add_br()
            # p.font.size = Pt(22)
            # p.font.color.rgb = RGBColor(255, 255, 255)  # White text color
            # p.font.name = 'Georgia'
            # p.space_after = Pt(14)  # Add space after each paragraph

            # Indent every second paragraph by 2 cm
            if i % 2 == 1:
                p.level = 1  # This will create an indent

            if i == len(self.verses) - 1:
                # At the last verse, append " - Halleluja" in italic
                if not self.verses[i][-1].endswith("Halleluja!"):
                    halleluja_run = p.add_run()
                    halleluja_run.font.italic = True
                    halleluja_run.text = " — Halleluja!"

        # Save the presentation
        self.prs.save(f'Psalm_{self.psalm_number:03d}.pptx')


def process_psalm(psalm: int, body_font_size, line_spacing_factor, space_after, wrap):
    try:
        prs = Presentation("Template.pptx")
        pw = PsalmWriter(psalm, prs, body_font_size, line_spacing_factor, space_after, wrap)
        pw.write_psalm()
    except Exception as e:
        print(f"Error while writing psalm {psalm}:\n")
        raise e


if __name__ == "__main__":
    standard_spacing_font_specific = 1.5  # Depending on font!
    spacing_factor = 1.2
    line_spacing_factor = standard_spacing_font_specific * spacing_factor

    func = partial(process_psalm, body_font_size=23, line_spacing_factor=line_spacing_factor, space_after=12, wrap=89)

    if PSALM is None:
        # Use process_map for multiprocessing with a progress bar
        results = process_map(func, range(1, 151), max_workers=cpu_count(), chunksize=1)
    else:
        prs = Presentation("Template.pptx")
        pw = PsalmWriter(PSALM, prs, body_font_size=23, line_spacing_factor=line_spacing_factor, space_after=12,
                         wrap=89)
        pw.write_psalm()
